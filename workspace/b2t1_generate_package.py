from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
import os, textwrap, json, zipfile, subprocess, shutil

OUT='/root/Hermes-Agent/workspace/B2T1_Navigating_Digital_Workspace'
os.makedirs(OUT, exist_ok=True)

meta={
 'topic_id':'B2T1',
 'topic_title':'Navigating the Digital Workspace',
 'duration':'50 minutes',
 'trades':'Mechanic Diesel and Mechanic Motor Vehicle',
 'audience':'ITI / NCVET learners, 10th or 12th pass, rural or first-time digital learners',
 'constraints':['Simple English','No Google or related product mentions','No vehicle brand names','Trade-linked examples throughout']
}

los=[
 'Explain what a digital workspace means in an automobile workshop.',
 'Identify common digital tools used in diesel and motor vehicle service work, such as job cards, service records, diagnostic reports, checklists, and parts lists.',
 'Read simple digital information such as an alert, fault code note, service history, or online instruction safely.',
 'Organise workshop files and records using clear names, folders, and dates.',
 'Follow safe behaviour while using digital devices: verify information, protect customer data, and ask a senior when unsure.'
]

use_cases=[
 {'name':'Use Case 1: Digital job card at the workshop', 'short':'A learner receives a job card on a tablet or computer. It shows vehicle complaint, customer note, service due items, and assigned work.', 'lesson':'Read the job card first, confirm the work, update status clearly, and do not skip safety checks.'},
 {'name':'Use Case 2: Diagnostic scan report for a diesel engine', 'short':'A scan tool gives a fault code and short message, for example sensor signal low or fuel pressure reading not normal.', 'lesson':'The digital message is a clue, not the final answer. Verify with visual inspection, meter readings, and instructor guidance.'},
 {'name':'Use Case 3: Parts and service record search', 'short':'A workshop assistant checks stock, part number, previous service date, and repair notes before starting work.', 'lesson':'Correct records save time, avoid wrong parts, and help the next technician understand what was done.'}
]

video_script_en='''
Imagine you are entering a modern automobile workshop for your first job. You see spanners, lifts, engines, batteries, and tools. You also see one more important tool: a digital workspace. This may be a computer, tablet, mobile app, diagnostic scanner, or a shared screen used by the workshop team.

A digital workspace helps a mechanic see the job card, customer complaint, service history, diagnostic report, parts list, safety checklist, and work status. It does not replace the mechanic. It helps the mechanic work in a more organised and safe way.

Let us take a simple example. A diesel vehicle comes to the workshop with low pickup. The job card says: customer reports low power on slope. The technician checks the service history and sees that the fuel filter was not changed in the last service. A scan report also shows an engine sensor warning. Now the technician has useful clues. But the technician should not blindly change parts. The correct step is to inspect, test, discuss with the senior, and then repair.

Another example is a motor vehicle service bay. A learner is asked to update a digital checklist. The checklist asks: engine oil checked, brake fluid checked, battery terminals checked, tyre condition checked, tools returned. By marking each item honestly, the learner helps the whole workshop avoid mistakes.

A good digital workspace has three habits. First, read carefully. Second, organise files with clear names, dates, and job numbers. Third, verify before acting. Digital information can be helpful, but it can also be incomplete or wrong. A safe mechanic uses digital information along with eyes, hands, measuring tools, and guidance from seniors.

In this topic, you will learn how to move around a digital workspace, read simple workshop information, organise records, and use digital tools safely. These skills will help you in class, in the lab, and later at work.
'''.strip()

video_script_mr='''
कल्पना करा की तुम्ही एका आधुनिक वाहन दुरुस्ती कार्यशाळेत पहिल्या दिवशी जाता. तिथे स्पॅनर, लिफ्ट, इंजिन, बॅटरी आणि साधने दिसतात. यासोबत अजून एक महत्त्वाचे साधन दिसते: डिजिटल कार्यक्षेत्र. हे संगणक, टॅबलेट, मोबाइल अॅप, डायग्नोस्टिक स्कॅनर किंवा कार्यशाळेतील सामायिक स्क्रीन असू शकते.

डिजिटल कार्यक्षेत्र मेकॅनिकला जॉब कार्ड, ग्राहकाची तक्रार, सर्व्हिस इतिहास, डायग्नोस्टिक रिपोर्ट, पार्ट्स यादी, सुरक्षा चेकलिस्ट आणि कामाची स्थिती पाहण्यास मदत करते. हे मेकॅनिकची जागा घेत नाही. ते मेकॅनिकला अधिक नीटनेटके आणि सुरक्षित काम करायला मदत करते.

एक सोपे उदाहरण पाहू. डिझेल वाहन कमी पिकअपच्या तक्रारीसह कार्यशाळेत येते. जॉब कार्डवर लिहिले आहे: चढावर वाहनाची शक्ती कमी वाटते. तंत्रज्ञ सर्व्हिस इतिहास पाहतो आणि मागील सर्व्हिसमध्ये इंधन फिल्टर बदललेले नाही हे दिसते. स्कॅन रिपोर्टमध्ये इंजिन सेन्सरची सूचना देखील दिसते. आता तंत्रज्ञाकडे उपयोगी माहिती आहे. पण भाग लगेच बदलणे योग्य नाही. योग्य पायरी म्हणजे तपासणी करणे, मोजमाप करणे, वरिष्ठांशी चर्चा करणे आणि मग दुरुस्ती करणे.

दुसरे उदाहरण मोटर वाहन सर्व्हिस बेचे आहे. शिकणाऱ्याला डिजिटल चेकलिस्ट अपडेट करायला सांगितले जाते. चेकलिस्ट विचारते: इंजिन ऑइल तपासले का, ब्रेक फ्लुइड तपासले का, बॅटरी टर्मिनल तपासले का, टायरची स्थिती पाहिली का, साधने परत ठेवली का. प्रत्येक बाब प्रामाणिकपणे चिन्हांकित केल्यास कार्यशाळेतील चुका कमी होतात.

चांगल्या डिजिटल कार्यक्षेत्रासाठी तीन सवयी महत्त्वाच्या आहेत. पहिले, काळजीपूर्वक वाचा. दुसरे, फाइल्सना स्पष्ट नाव, तारीख आणि जॉब नंबर द्या. तिसरे, कृती करण्यापूर्वी पडताळणी करा. डिजिटल माहिती उपयोगी असते, पण ती कधी कधी अपूर्ण किंवा चुकीची असू शकते. सुरक्षित मेकॅनिक डिजिटल माहितीबरोबर स्वतःचे निरीक्षण, हाताने तपासणी, मोजमाप साधने आणि वरिष्ठांचे मार्गदर्शन वापरतो.

या विषयात तुम्ही डिजिटल कार्यक्षेत्रात कसे फिरायचे, साधी कार्यशाळा माहिती कशी वाचायची, नोंदी कशा व्यवस्थित ठेवायच्या आणि डिजिटल साधने सुरक्षितपणे कशी वापरायची हे शिकाल. ही कौशल्ये वर्गात, लॅबमध्ये आणि नोकरीत उपयोगी पडतील.
'''.strip()

rm1_en='''
# Reading Material 1: What is a Digital Workspace in an Automobile Workshop?

## 1. A workshop is no longer only physical
In an automobile workshop, learners usually think of tools, engines, wheels, batteries, fuel systems, brakes, and measuring instruments. These are still very important. But many workshops also use digital tools. A digital workspace is the place where work information is stored, read, updated, and shared using a device.

A digital workspace may include:
- A digital job card
- Service history of a vehicle
- A diagnostic scan report
- A digital checklist
- A parts stock record
- Photos of damaged parts
- Safety instructions
- Work status, such as pending, in progress, or completed

For Mechanic Diesel and Mechanic Motor Vehicle learners, this is useful because modern maintenance work needs both hand skills and information skills.

## 2. Why it matters
A mechanic may do the correct repair only when the correct information is understood. If the job card is not read properly, the wrong complaint may be checked. If the service record is ignored, an old problem may be missed. If the diagnostic report is followed blindly, a good part may be replaced unnecessarily.

Digital tools help in four ways:
1. They save time by keeping information in one place.
2. They reduce mistakes by using checklists and records.
3. They help the team communicate clearly.
4. They support safety by showing warnings and instructions.

## 3. Use Case 1: Digital job card
A diesel vehicle comes to the workshop. The digital job card says: customer reports hard starting in the morning. The job card also shows last service date, fuel filter status, battery check status, and assigned technician.

A good learner should:
- Read the complaint fully.
- Check the service history.
- Note important observations.
- Ask the instructor or senior before taking major action.
- Update the job card only with true information.

This simple habit prevents confusion. The next technician can understand what was checked and what is still pending.

## 4. Common areas in a digital workspace
Most digital workspaces have similar areas:
- Home screen: shows main options.
- Search: helps find a job number, vehicle number, part, or service record.
- Folders: store documents, photos, and reports.
- Forms: collect information in a fixed format.
- Upload button: adds a photo or document.
- Save button: stores your work.
- Status field: shows whether the job is open, waiting, or completed.

## 5. Use this in class or lab
During a practical session, make a paper or digital job card for a simple task, such as checking battery terminals or inspecting air filter condition. Ask yourself:
- What is the complaint?
- What safety check is needed first?
- What did I observe?
- What action did I take?
- What should the next person know?

## 6. Trust, but verify
Digital information is helpful, but it is not always complete. A scan report may show a sensor warning, but the real cause may be loose wiring, weak battery, damaged connector, or poor maintenance. A safe mechanic reads the digital clue and then verifies it with inspection and testing.

Remember: Digital tools guide the mechanic. They do not replace the mechanic’s thinking.
'''.strip()

rm2_en='''
# Reading Material 2: Reading, Organising, and Using Digital Information Safely

## 1. Reading digital information
A digital workspace may show short messages. Read them slowly. Do not rush. Look for five things:
1. What is the problem or task?
2. Which system is involved? Engine, fuel, brake, battery, cooling, or electrical?
3. What is the date and job number?
4. What action has already been taken?
5. What needs verification before repair?

## 2. Use Case 2: Diagnostic scan report
A diagnostic scanner gives a message related to an engine sensor. This does not mean the sensor must be replaced immediately. It means the system has found an unusual reading.

Correct method:
- Read the message.
- Write down or save the code/report.
- Check related wiring, connector, fuse, sensor fitting, and visible damage.
- Compare with service instructions given by the instructor or workshop.
- Take help from a senior before replacing parts.

This method saves money and avoids wrong repair.

## 3. Organising files and photos
A workshop may keep photos of damaged parts, service reports, and checklist records. If files are named badly, they become difficult to find.

Good file name example:
`Job1245_2026-05-15_BatteryTerminal_Photo1`

Poor file name example:
`image123_new_final_last`

Use clear names with:
- Job number
- Date
- System or part name
- Type of record

## 4. Use Case 3: Parts and service record search
A vehicle comes for service. The learner checks the service record and sees that brake fluid was checked last time, but air filter replacement was pending. The parts list shows the correct part is available in stock. This helps the workshop plan work without delay.

If records are wrong, the wrong part may be selected or the same work may be repeated. Correct digital records save time and reduce waste.

## 5. Safe digital habits
Follow these habits in any digital workspace:
- Do not share customer phone numbers, addresses, or vehicle details outside work.
- Do not change records unless you are asked to update them.
- Do not mark a checklist item as completed if you have not checked it.
- Do not blindly follow a digital alert. Verify with practical checks.
- Ask for help if you do not understand a message.
- Keep the device clean and away from oil, water, heat, and moving parts.

## 6. Smooth workflow for a learner
A simple workflow is:
1. Open the job card.
2. Read the complaint and task.
3. Check safety instructions.
4. View service history or report.
5. Do the practical inspection.
6. Update observations clearly.
7. Ask the senior to verify.
8. Save and close the job.

## 7. Did you know?
Many workshops now use digital records to reduce repeat complaints. When the next technician can read what happened earlier, the customer gets faster and better service.

Digital skill is now part of workshop discipline, just like tool discipline and safety discipline.
'''.strip()

summary_en='''
# Topic Summary: Navigating the Digital Workspace

In this topic, you learned that a digital workspace is a place where workshop information is stored, read, updated, and shared using a device. For Mechanic Diesel and Mechanic Motor Vehicle learners, it may include job cards, service history, scan reports, checklists, parts records, photos, and safety instructions.

Main points:
- A digital workspace helps the mechanic work in an organised way.
- A digital job card tells the complaint, task, service history, and status.
- A diagnostic report gives clues, but the mechanic must verify before repair.
- Good file names and correct records help the whole workshop team.
- Safety and privacy are important. Do not share customer details or mark false checklist items.
- Digital tools support the mechanic. They do not replace practical skill, observation, and senior guidance.

Remember the three habits: Read carefully, organise clearly, and verify before action.
'''.strip()

assignment_en='''
# Assignment: Find the Digital Workspace Around You

Time: 5 minutes

Imagine you are in your ITI lab or a small automobile workshop. Choose one task:
- Battery terminal inspection
- Air filter inspection
- Engine oil level check
- Brake fluid level check
- Diesel fuel filter check

Create a simple digital or paper job card with these fields:
1. Job number
2. Date
3. Task name
4. Customer complaint or training task
5. Safety check before work
6. Observation
7. Action taken
8. What needs senior verification?

Submission: Write 5 to 7 lines explaining how this record will help the next technician.
'''.strip()

assessment=[
('What is a digital workspace in a workshop?', ['A place to sleep','A place/device system where job information is read, stored, updated, and shared','Only a toolbox','Only a fuel tank'],2),
('A digital job card says “hard starting in morning”. What should the learner do first?', ['Ignore it','Read the full complaint and service history','Replace all parts','Close the job'],2),
('A scan report gives an engine sensor warning. What is the safest action?', ['Replace the sensor immediately','Use it as a clue and verify with inspection/testing','Delete the report','Tell the customer there is no problem'],2),
('Which file name is clearer?', ['image123_final','Job1245_2026-05-15_BatteryTerminal_Photo1','new photo','last last file'],2),
('Why should checklist items be marked honestly?', ['To finish fast','To avoid safety mistakes and false records','To hide incomplete work','To confuse the team'],2),
('Which information should not be shared outside work?', ['Customer phone number and vehicle details','Tool name','General safety rule','Class topic name'],1),
('Digital tools replace mechanic thinking. Is this correct?', ['Yes','No, they support the mechanic but practical verification is needed','Only on Mondays','Only for diesel vehicles'],2),
('What should a learner do if they do not understand a digital message?', ['Guess the repair','Ask instructor or senior and verify','Delete the message','Ignore safety steps'],2),
('Which habit helps find records quickly?', ['Clear folders and file names','Random names','No dates','Saving everything on one screen'],1),
('The correct workflow starts with:', ['Repair first and read later','Open and read the job card','Mark completed before checking','Share customer details'],2)
]

# manual Marathi translations
mr = {
 'title':'B2T1 - डिजिटल कार्यक्षेत्रात मार्गदर्शन',
 'rm1': rm1_en.replace('Reading Material 1: What is a Digital Workspace in an Automobile Workshop?','वाचन साहित्य 1: वाहन कार्यशाळेत डिजिटल कार्यक्षेत्र म्हणजे काय?')
    .replace('digital workspace','डिजिटल कार्यक्षेत्र').replace('Digital workspace','डिजिटल कार्यक्षेत्र').replace('job card','जॉब कार्ड').replace('Job card','जॉब कार्ड'),
}
# More natural Marathi docs below, not raw replace only
rm1_mr='''
# वाचन साहित्य 1: वाहन कार्यशाळेत डिजिटल कार्यक्षेत्र म्हणजे काय?

## 1. कार्यशाळा फक्त साधनांची जागा राहिलेली नाही
वाहन कार्यशाळेत साधने, इंजिन, चाके, बॅटरी, इंधन प्रणाली, ब्रेक आणि मोजमाप साधने महत्त्वाची असतात. पण आज अनेक कार्यशाळांमध्ये डिजिटल साधनेही वापरली जातात. डिजिटल कार्यक्षेत्र म्हणजे अशी जागा किंवा प्रणाली जिथे कामाची माहिती उपकरणावर वाचली, साठवली, अपडेट केली आणि टीमसोबत शेअर केली जाते.

यात जॉब कार्ड, सर्व्हिस इतिहास, डायग्नोस्टिक रिपोर्ट, डिजिटल चेकलिस्ट, पार्ट्स स्टॉक रेकॉर्ड, खराब भागांचे फोटो, सुरक्षा सूचना आणि कामाची स्थिती यांचा समावेश असू शकतो.

## 2. हे का महत्त्वाचे आहे?
योग्य माहिती समजली तरच योग्य दुरुस्ती करता येते. जॉब कार्ड नीट वाचले नाही तर चुकीची तक्रार तपासली जाऊ शकते. सर्व्हिस रेकॉर्ड दुर्लक्ष केले तर जुनी समस्या राहू शकते. डायग्नोस्टिक रिपोर्ट आंधळेपणाने वापरला तर अनावश्यक भाग बदलला जाऊ शकतो.

डिजिटल साधने वेळ वाचवतात, चुका कमी करतात, टीमला स्पष्ट संवादात मदत करतात आणि सुरक्षा सूचना दाखवतात.

## 3. वापर उदाहरण 1: डिजिटल जॉब कार्ड
एक डिझेल वाहन कार्यशाळेत येते. जॉब कार्डवर लिहिले आहे: सकाळी वाहन सुरू होण्यास अडचण. त्यात मागील सर्व्हिस तारीख, इंधन फिल्टरची स्थिती, बॅटरी तपासणी आणि काम दिलेल्या तंत्रज्ञाचे नाव दिसते.

चांगल्या शिकणाऱ्याने तक्रार पूर्ण वाचावी, सर्व्हिस इतिहास तपासावा, निरीक्षण लिहावे, मोठी कृती करण्यापूर्वी वरिष्ठांना विचारावे आणि खरी माहितीच अपडेट करावी.

## 4. डिजिटल कार्यक्षेत्रातील सामान्य भाग
होम स्क्रीन, शोध, फोल्डर, फॉर्म, अपलोड, सेव्ह आणि स्थिती हे भाग सामान्य असतात. यांचा वापर करून जॉब नंबर, वाहन नोंद, भाग, फोटो आणि रिपोर्ट शोधता येतात.

## 5. वर्ग किंवा लॅबमध्ये वापरा
बॅटरी टर्मिनल तपासणी किंवा एअर फिल्टर तपासणीसाठी साधे जॉब कार्ड तयार करा. तक्रार काय आहे, आधी कोणती सुरक्षा तपासणी करायची, काय निरीक्षण झाले, काय कृती केली आणि पुढच्या व्यक्तीला काय माहिती हवी हे लिहा.

## 6. विश्वास ठेवा, पण पडताळणी करा
डिजिटल माहिती उपयोगी असते, पण ती नेहमी पूर्ण असेलच असे नाही. स्कॅन रिपोर्ट सेन्सर सूचना दाखवू शकतो, पण खरे कारण सैल वायर, कमजोर बॅटरी, खराब कनेक्टर किंवा कमी देखभाल असू शकते. सुरक्षित मेकॅनिक डिजिटल सूचना वाचतो आणि मग प्रत्यक्ष तपासणी व चाचणी करतो.
'''.strip()
rm2_mr='''
# वाचन साहित्य 2: डिजिटल माहिती वाचणे, व्यवस्थित ठेवणे आणि सुरक्षित वापरणे

## 1. डिजिटल माहिती वाचणे
डिजिटल कार्यक्षेत्रात लहान संदेश दिसू शकतात. ते हळू वाचा. पाच गोष्टी पहा: समस्या काय आहे, कोणती प्रणाली संबंधित आहे, तारीख आणि जॉब नंबर काय आहे, आधी कोणती कृती झाली आहे आणि दुरुस्तीपूर्वी काय पडताळायचे आहे.

## 2. वापर उदाहरण 2: डायग्नोस्टिक स्कॅन रिपोर्ट
स्कॅनर इंजिन सेन्सरशी संबंधित संदेश दाखवतो. याचा अर्थ सेन्सर लगेच बदलायचा असा नाही. याचा अर्थ प्रणालीला असामान्य रीडिंग दिसले आहे.

योग्य पद्धत म्हणजे संदेश वाचा, कोड किंवा रिपोर्ट सेव्ह करा, वायरिंग, कनेक्टर, फ्यूज, सेन्सर फिटिंग आणि दिसणारे नुकसान तपासा, सूचना वरिष्ठ किंवा प्रशिक्षकासोबत पडताळा आणि मगच भाग बदलण्याचा निर्णय घ्या.

## 3. फाइल्स आणि फोटो व्यवस्थित ठेवणे
खराब भागांचे फोटो, सर्व्हिस रिपोर्ट आणि चेकलिस्ट रेकॉर्ड ठेवताना स्पष्ट नाव द्या. चांगले नाव: `Job1245_2026-05-15_BatteryTerminal_Photo1`. यात जॉब नंबर, तारीख, भागाचे नाव आणि रेकॉर्डचा प्रकार दिसतो.

## 4. वापर उदाहरण 3: पार्ट्स आणि सर्व्हिस रेकॉर्ड शोध
वाहन सर्व्हिससाठी येते. शिकणारा सर्व्हिस रेकॉर्ड पाहतो आणि ब्रेक फ्लुइड मागील वेळी तपासलेले आहे, पण एअर फिल्टर बदलणे बाकी आहे हे दिसते. पार्ट्स यादीत योग्य भाग उपलब्ध आहे. यामुळे काम वेळेवर सुरू होते.

## 5. सुरक्षित डिजिटल सवयी
ग्राहकाचा फोन नंबर किंवा वाहन माहिती बाहेर शेअर करू नका. विचारल्याशिवाय रेकॉर्ड बदलू नका. तपासणी न करता चेकलिस्ट पूर्ण म्हणून चिन्हांकित करू नका. डिजिटल सूचना आंधळेपणाने मानू नका. समजले नाही तर वरिष्ठांना विचारा. उपकरण तेल, पाणी, उष्णता आणि हलणाऱ्या भागांपासून दूर ठेवा.

## 6. शिकणाऱ्यासाठी सोपा प्रवाह
जॉब कार्ड उघडा, तक्रार वाचा, सुरक्षा सूचना पहा, सर्व्हिस इतिहास किंवा रिपोर्ट पहा, प्रत्यक्ष तपासणी करा, निरीक्षण स्पष्ट लिहा, वरिष्ठांकडून पडताळा आणि जॉब सेव्ह करून बंद करा.

डिजिटल कौशल्य हे साधन शिस्त आणि सुरक्षा शिस्तीसारखेच कार्यशाळेतील महत्त्वाचे कौशल्य आहे.
'''.strip()
summary_mr='''
# विषय सारांश: डिजिटल कार्यक्षेत्रात मार्गदर्शन

या विषयात तुम्ही शिकलात की डिजिटल कार्यक्षेत्रात कार्यशाळेची माहिती उपकरणावर वाचली, साठवली, अपडेट केली आणि शेअर केली जाते. Mechanic Diesel आणि Mechanic Motor Vehicle शिकणाऱ्यांसाठी यात जॉब कार्ड, सर्व्हिस इतिहास, स्कॅन रिपोर्ट, चेकलिस्ट, पार्ट्स रेकॉर्ड, फोटो आणि सुरक्षा सूचना असू शकतात.

मुख्य मुद्दे:
- डिजिटल कार्यक्षेत्र मेकॅनिकला नीटनेटके काम करायला मदत करते.
- जॉब कार्ड तक्रार, काम, सर्व्हिस इतिहास आणि स्थिती दाखवते.
- डायग्नोस्टिक रिपोर्ट सूचना देतो, पण दुरुस्तीपूर्वी पडताळणी आवश्यक आहे.
- स्पष्ट फाइल नावे आणि बरोबर रेकॉर्ड संपूर्ण टीमला मदत करतात.
- ग्राहकाची माहिती सुरक्षित ठेवणे महत्त्वाचे आहे.
- डिजिटल साधने मदत करतात, पण प्रत्यक्ष कौशल्य आणि वरिष्ठांचे मार्गदर्शन आवश्यक आहे.

तीन सवयी लक्षात ठेवा: काळजीपूर्वक वाचा, स्पष्टपणे व्यवस्थित ठेवा आणि कृतीपूर्वी पडताळा.
'''.strip()
assignment_mr='''
# असाइनमेंट: तुमच्या आसपासचे डिजिटल कार्यक्षेत्र शोधा

वेळ: 5 मिनिटे

तुम्ही ITI लॅबमध्ये किंवा छोट्या वाहन कार्यशाळेत आहात असे समजा. एक काम निवडा: बॅटरी टर्मिनल तपासणी, एअर फिल्टर तपासणी, इंजिन ऑइल लेव्हल तपासणी, ब्रेक फ्लुइड लेव्हल तपासणी किंवा डिझेल फ्युएल फिल्टर तपासणी.

या फील्डसह साधे जॉब कार्ड तयार करा: जॉब नंबर, तारीख, कामाचे नाव, तक्रार किंवा प्रशिक्षण काम, कामापूर्वी सुरक्षा तपासणी, निरीक्षण, केलेली कृती आणि वरिष्ठांनी काय पडताळायचे.

सबमिशन: हा रेकॉर्ड पुढच्या तंत्रज्ञाला कसा मदत करेल यावर 5 ते 7 ओळी लिहा.
'''.strip()

# docx helpers
def add_title(doc, title):
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(title); r.bold=True; r.font.size=Pt(18); r.font.color.rgb=RGBColor(20,80,70)

def add_md(doc, md):
    for line in md.split('\n'):
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.strip()=='' :
            doc.add_paragraph('')
        else:
            doc.add_paragraph(line)

def make_doc(name, title, sections):
    doc=Document(); add_title(doc,title)
    for heading, content in sections:
        doc.add_heading(heading, level=1)
        add_md(doc, content)
    path=os.path.join(OUT,name); doc.save(path); return path

make_doc('B2T1_English_Content_Package.docx','B2T1 - Navigating the Digital Workspace',[
 ('Topic Metadata', json.dumps(meta, indent=2)),
 ('Learning Outcomes', '\n'.join('- '+x for x in los)),
 ('Integrated Use Cases', '\n'.join([f"- {u['name']}: {u['short']} Lesson: {u['lesson']}" for u in use_cases])),
 ('Topic Introduction - 5 mins','''Welcome to B2T1 - Navigating the Digital Workspace. In this topic, learners connect familiar workshop work with digital records, reports, checklists, and safe information use.'''),
 ('Learning Video Script - 3 to 4 mins', video_script_en),
 ('Reading Material 1 - 12 to 15 mins', rm1_en),
 ('Reading Material 2 - 12 to 15 mins', rm2_en),
 ('Topic Summary - 5 mins', summary_en),
 ('Assignment - 5 mins', assignment_en)
])
make_doc('B2T1_Marathi_Content_Package.docx','B2T1 - डिजिटल कार्यक्षेत्रात मार्गदर्शन',[
 ('विषय माहिती', 'कालावधी: 50 मिनिटे\nट्रेड: Mechanic Diesel आणि Mechanic Motor Vehicle\nभाषा: सोपी मराठी'),
 ('शिकण्याचे परिणाम', '\n'.join('- '+x for x in ['डिजिटल कार्यक्षेत्र म्हणजे काय हे सांगणे','जॉब कार्ड, सर्व्हिस इतिहास, स्कॅन रिपोर्ट आणि चेकलिस्ट ओळखणे','डिजिटल सूचना वाचून पडताळणीची गरज समजणे','फाइल्स आणि रेकॉर्ड स्पष्टपणे व्यवस्थित ठेवणे','ग्राहक माहिती आणि डिजिटल सुरक्षिततेच्या सवयी पाळणे'])),
 ('व्हिडिओ स्क्रिप्ट - 3 ते 4 मिनिटे', video_script_mr),
 ('वाचन साहित्य 1 - 12 ते 15 मिनिटे', rm1_mr),
 ('वाचन साहित्य 2 - 12 ते 15 मिनिटे', rm2_mr),
 ('विषय सारांश - 5 मिनिटे', summary_mr),
 ('असाइनमेंट - 5 मिनिटे', assignment_mr)
])

# Assessment xlsx
wb=Workbook(); ws=wb.active; ws.title='Assessment_EN'
headers=['Question No','Question Text','Option Count','Correct Option','Option 1','Option 2','Option 3','Option 4','Module no']
ws.append(headers)
for i,(q,opts,correct) in enumerate(assessment,1): ws.append([i,q,4,correct]+opts+['B2T1'])
ws2=wb.create_sheet('Blueprint')
ws2.append(['Topic','Duration','Level','Notes'])
ws2.append(['B2T1 Navigating the Digital Workspace','10 mins','Knowledge + Comprehension','10 MCQs, simple language, trade-linked'])
wb.save(os.path.join(OUT,'B2T1_Assessment_EN.xlsx'))

# markdown files
for name, content in {
 'B2T1_Video_Script_EN.md': video_script_en,
 'B2T1_Video_Script_MR.md': video_script_mr,
 'B2T1_RM1_EN.md': rm1_en,
 'B2T1_RM2_EN.md': rm2_en,
 'B2T1_RM1_MR.md': rm1_mr,
 'B2T1_RM2_MR.md': rm2_mr,
 'B2T1_Summary_EN.md': summary_en,
 'B2T1_Summary_MR.md': summary_mr,
 'B2T1_Assignment_EN.md': assignment_en,
 'B2T1_Assignment_MR.md': assignment_mr,
}.items():
    open(os.path.join(OUT,name),'w',encoding='utf-8').write(content+'\n')

# PPTX
def add_slide(prs, title, body, num, lang='EN'):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=PRGB(246,248,242)
    accent=slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(0.22), PInches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb=PRGB(27,94,76); accent.line.fill.background()
    tx=slide.shapes.add_textbox(PInches(0.55),PInches(0.35),PInches(8.4),PInches(0.8)).text_frame
    tx.text=title; tx.paragraphs[0].font.bold=True; tx.paragraphs[0].font.size=PPt(30); tx.paragraphs[0].font.color.rgb=PRGB(31,55,48)
    box=slide.shapes.add_shape(5,PInches(0.65),PInches(1.35),PInches(7.4),PInches(4.8))
    box.fill.solid(); box.fill.fore_color.rgb=PRGB(255,255,255); box.line.color.rgb=PRGB(210,220,210)
    tf=box.text_frame; tf.margin_left=PInches(0.25); tf.margin_right=PInches(0.2); tf.margin_top=PInches(0.18)
    for j,line in enumerate(body):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph(); p.text=line; p.font.size=PPt(18 if len(line)<90 else 15); p.font.color.rgb=PRGB(40,55,50); p.space_after=PPt(8)
    icon=slide.shapes.add_shape(9,PInches(8.45),PInches(5.45),PInches(1.0),PInches(1.0))
    icon.fill.solid(); icon.fill.fore_color.rgb=PRGB(244,171,66); icon.line.fill.background()
    it=icon.text_frame; it.text=str(num); it.paragraphs[0].font.bold=True; it.paragraphs[0].font.size=PPt(28); it.paragraphs[0].alignment=PP_ALIGN.CENTER

def make_pptx(path, lang='EN'):
    prs=Presentation(); prs.slide_width=PInches(10); prs.slide_height=PInches(7.5)
    if lang=='EN':
        slides=[
            ('Navigating the Digital Workspace',['For Mechanic Diesel and Mechanic Motor Vehicle learners','50-minute topic','Simple rule: read, organise, verify']),
            ('Why this matters',['Modern workshops use both hand tools and digital tools.','Digital records help reduce mistakes and save time.','A mechanic still needs practical skill and judgement.']),
            ('What you may see',['Digital job card','Service history','Diagnostic scan report','Checklist and parts record','Photos and safety instructions']),
            ('Use Case 1: Digital job card',['Read customer complaint fully.','Check service history before starting.','Update true observations only.']),
            ('Use Case 2: Scan report',['A fault message is a clue.','Verify wiring, connector, fuse, fitting, and readings.','Ask a senior before replacing parts.']),
            ('Use Case 3: Parts and records',['Search previous service date and pending work.','Use correct part and stock record.','Good records help the next technician.']),
            ('Safe digital habits',['Protect customer details.','Do not mark work complete without checking.','Keep devices away from oil, water, and heat.']),
            ('Learner workflow',['Open job card → read task → check safety','Inspect practically → update observations','Get senior verification → save and close']),
            ('Remember',['Digital tools support the mechanic.','They do not replace observation and testing.','Read carefully, organise clearly, verify before action.'])]
    else:
        slides=[
            ('डिजिटल कार्यक्षेत्रात मार्गदर्शन',['Mechanic Diesel आणि Mechanic Motor Vehicle शिकणाऱ्यांसाठी','50 मिनिटांचा विषय','नियम: वाचा, व्यवस्थित ठेवा, पडताळा']),
            ('हे का महत्त्वाचे आहे?',['आधुनिक कार्यशाळेत हाताची साधने आणि डिजिटल साधने दोन्ही वापरली जातात.','डिजिटल रेकॉर्ड चुका कमी करतात.','मेकॅनिकचे कौशल्य आणि विचार महत्त्वाचेच राहतात.']),
            ('तुम्हाला काय दिसू शकते?',['डिजिटल जॉब कार्ड','सर्व्हिस इतिहास','डायग्नोस्टिक स्कॅन रिपोर्ट','चेकलिस्ट आणि पार्ट्स रेकॉर्ड','फोटो आणि सुरक्षा सूचना']),
            ('वापर उदाहरण 1: जॉब कार्ड',['ग्राहकाची तक्रार पूर्ण वाचा.','काम सुरू करण्यापूर्वी सर्व्हिस इतिहास तपासा.','खरी निरीक्षणेच अपडेट करा.']),
            ('वापर उदाहरण 2: स्कॅन रिपोर्ट',['फॉल्ट संदेश हा फक्त संकेत असतो.','वायरिंग, कनेक्टर, फ्यूज आणि रीडिंग पडताळा.','भाग बदलण्यापूर्वी वरिष्ठांना विचारा.']),
            ('वापर उदाहरण 3: पार्ट्स आणि रेकॉर्ड',['मागील सर्व्हिस तारीख आणि बाकी काम शोधा.','योग्य भाग आणि स्टॉक रेकॉर्ड वापरा.','चांगले रेकॉर्ड पुढच्या तंत्रज्ञाला मदत करतात.']),
            ('सुरक्षित डिजिटल सवयी',['ग्राहक माहिती सुरक्षित ठेवा.','तपासणी न करता पूर्ण चिन्हांकित करू नका.','उपकरण तेल, पाणी आणि उष्णतेपासून दूर ठेवा.']),
            ('शिकणाऱ्याचा प्रवाह',['जॉब कार्ड उघडा → काम वाचा → सुरक्षा तपासा','प्रत्यक्ष तपासणी करा → निरीक्षण लिहा','वरिष्ठ पडताळणी → सेव्ह आणि बंद']),
            ('लक्षात ठेवा',['डिजिटल साधने मेकॅनिकला मदत करतात.','ती निरीक्षण आणि चाचणीची जागा घेत नाहीत.','काळजीपूर्वक वाचा, व्यवस्थित ठेवा, पडताळा.'])]
    for i,(t,b) in enumerate(slides,1): add_slide(prs,t,b,i,lang)
    prs.save(path)
make_pptx(os.path.join(OUT,'B2T1_Learning_Video_Deck_EN.pptx'),'EN')
make_pptx(os.path.join(OUT,'B2T1_Learning_Video_Deck_MR.pptx'),'MR')

# create slide images with PIL for video draft
SLIDE_DIR=os.path.join(OUT,'video_slides_en'); os.makedirs(SLIDE_DIR, exist_ok=True)
try: font_title=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',52); font_body=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',34)
except: font_title=font_body=None
slides_simple=[('Navigating the Digital Workspace','Modern workshops use job cards, service records, scan reports, checklists and parts records.'),('Digital job card','Read the complaint, check service history, update true observations.'),('Scan report','A digital warning is a clue. Verify with inspection and testing.'),('Organise records','Use clear names with job number, date and part/system name.'),('Safe habits','Protect customer information. Do not mark checks as complete unless done.'),('Remember','Read carefully. Organise clearly. Verify before action.')]
for i,(t,b) in enumerate(slides_simple,1):
    img=Image.new('RGB',(1280,720),(246,248,242)); d=ImageDraw.Draw(img)
    d.rectangle([0,0,130,720], fill=(27,94,76)); d.ellipse([1060,500,1190,630], fill=(244,171,66))
    d.text((170,80),t,fill=(31,55,48),font=font_title)
    y=210
    for line in textwrap.wrap(b, width=48):
        d.text((180,y),line,fill=(40,55,50),font=font_body); y+=50
    d.text((1100,535),str(i),fill=(20,20,20),font=font_title)
    img.save(os.path.join(SLIDE_DIR,f'slide_{i:02d}.png'))
# silent draft video
concat=os.path.join(OUT,'video_slides_en','concat.txt')
with open(concat,'w') as f:
    for i in range(1,len(slides_simple)+1):
        f.write(f"file 'slide_{i:02d}.png'\nduration 25\n")
    f.write(f"file 'slide_{len(slides_simple):02d}.png'\n")
video=os.path.join(OUT,'B2T1_Learning_Video_EN_silent_draft.mp4')
subprocess.run(['ffmpeg','-y','-f','concat','-safe','0','-i',concat,'-vf','fps=25,format=yuv420p','-t','150',video],stdout=subprocess.DEVNULL,stderr=subprocess.DEVNULL)

# production notes
notes=f'''# Production Notes and Blockers\n\nCompleted locally:\n- English and Marathi content packages.\n- English and Marathi PPTX decks.\n- Assessment XLSX.\n- English silent draft video from slide visuals.\n\nBlocked / needs credentials or accessible assets:\n- Drive templates folder requires sign-in in this environment, so exact templates could not be downloaded. Created template-ready files instead.\n- Google Workspace OAuth is not authenticated here.\n- ELEVENLABS_API_KEY is not present, so ElevenLabs audio could not be generated from this server.\n- SARVAM_API_KEY is not present, so Sarvam Marathi audio could not be generated from this server.\n- Remotion is not installed in this project. A Remotion build can be added after templates/audio are available.\n\nRecommended next step:\n1. Put these contents into the client templates from the Drive folder.\n2. Generate English voiceover from `B2T1_Video_Script_EN.md` using ElevenLabs.\n3. Generate Marathi voiceover from `B2T1_Video_Script_MR.md` using Sarvam.\n4. Use the PPTX decks as source slides for Remotion animation/video sync.\n'''
open(os.path.join(OUT,'PRODUCTION_NOTES.md'),'w',encoding='utf-8').write(notes)

# zip
zip_path='/root/Hermes-Agent/workspace/B2T1_Navigating_Digital_Workspace_package.zip'
with zipfile.ZipFile(zip_path,'w',zipfile.ZIP_DEFLATED) as z:
    for root,dirs,files in os.walk(OUT):
        for fn in files:
            p=os.path.join(root,fn); z.write(p, os.path.relpath(p, os.path.dirname(OUT)))
print(zip_path)
print('files', len(os.listdir(OUT)))
