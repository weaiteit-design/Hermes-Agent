import os, textwrap, zipfile, math, re
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle, PageBreak, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

BASE='/root/Hermes-Agent/workspace/B2T1_Navigating_Digital_Workspace'
OUT=os.path.join(BASE,'formatted_final')
IMG_DIR=os.path.join(OUT,'images')
os.makedirs(IMG_DIR, exist_ok=True)

ASSET_DIR=os.path.join(BASE,'extracted_assets')
HEADER=os.path.join(ASSET_DIR,'doc_d37e65ca46e5_4._EW_E&W_B1T1_RM_A2.pdf_p1_img1.png')
FOOTER=os.path.join(ASSET_DIR,'doc_d37e65ca46e5_4._EW_E&W_B1T1_RM_A2.pdf_p1_img2.png')

# Register font
FONT='/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'
FONT_B='/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'
pdfmetrics.registerFont(TTFont('DejaVu', FONT))
pdfmetrics.registerFont(TTFont('DejaVu-Bold', FONT_B))

def load_text(name):
    return open(os.path.join(BASE,name),encoding='utf-8').read().strip()

rm1=load_text('B2T1_RM1_EN.md')
rm2=load_text('B2T1_RM2_EN.md')
summary=load_text('B2T1_Summary_EN.md')
assignment=load_text('B2T1_Assignment_EN.md')
video_script=load_text('B2T1_Video_Script_EN.md')

# Simple relevant generated illustrations using local drawing when image API unavailable.
def make_illustration(path, title, scene):
    W,H=1400,760
    im=Image.new('RGB',(W,H),(245,248,244))
    d=ImageDraw.Draw(im)
    title_f=ImageFont.truetype(FONT_B,54)
    body_f=ImageFont.truetype(FONT,30)
    small_f=ImageFont.truetype(FONT,24)
    # background workshop
    d.rectangle([0,0,W,90], fill=(28,92,76))
    d.text((48,20),title,fill='white',font=title_f)
    d.rectangle([0,650,W,760], fill=(222,229,219))
    d.line([0,650,W,650], fill=(190,198,188), width=4)
    # generic vehicle silhouette
    d.rounded_rectangle([690,270,1240,530], radius=42, fill=(210,222,226), outline=(80,95,100), width=5)
    d.rectangle([790,220,1110,345], fill=(190,207,214), outline=(80,95,100), width=5)
    d.ellipse([770,500,900,630], fill=(40,45,48)); d.ellipse([1030,500,1160,630], fill=(40,45,48))
    d.ellipse([805,535,865,595], fill=(175,183,185)); d.ellipse([1065,535,1125,595], fill=(175,183,185))
    # learner body
    d.ellipse([170,195,260,285], fill=(141,93,55))
    d.rectangle([195,285,235,350], fill=(141,93,55))
    d.rounded_rectangle([115,340,335,595], radius=35, fill=(244,171,66), outline=(90,80,60), width=4)
    d.line([150,390,75,500], fill=(141,93,55), width=24)
    d.line([305,390,425,500], fill=(141,93,55), width=24)
    d.rectangle([60,490,170,540], fill=(60,65,70))
    # tablet/scanner
    d.rounded_rectangle([385,390,640,565], radius=20, fill=(36,48,54), outline=(20,20,20), width=5)
    d.rounded_rectangle([410,420,615,535], radius=10, fill=(235,248,250))
    # screen contents based on scene
    if scene=='jobcard':
        for i,txt in enumerate(['JOB CARD','Complaint','Service history','Safety checklist']):
            y=435+i*26
            d.rectangle([425,y,595,y+16], fill=(230,238,234))
            d.text((430,y-5),txt,fill=(28,92,76),font=small_f if i==0 else ImageFont.truetype(FONT,16))
        notes=['Read complaint first','Check service history','Update true observations']
    elif scene=='diagnostic':
        d.line([640,490,760,410], fill=(45,45,45), width=8)
        for i,txt in enumerate(['SCAN REPORT','Fault code clue','Verify wiring','Test before repair']):
            y=435+i*26
            d.rectangle([425,y,595,y+16], fill=(230,238,234))
            d.text((430,y-5),txt,fill=(28,92,76),font=small_f if i==0 else ImageFont.truetype(FONT,16))
        d.rectangle([825,380,1045,460], fill=(70,78,82)); d.text((850,402),'ENGINE', fill='white', font=body_f)
        notes=['Report is a clue','Inspect connector and wiring','Ask senior before replacing parts']
    else:
        d.rectangle([760,175,1240,255], fill=(180,135,75));
        for x in range(785,1210,100): d.rectangle([x,185,x+65,240], fill=(225,220,190), outline=(130,100,65))
        for i,txt in enumerate(['PARTS RECORD','Stock status','Previous service','Pending work']):
            y=435+i*26
            d.rectangle([425,y,595,y+16], fill=(230,238,234))
            d.text((430,y-5),txt,fill=(28,92,76),font=small_f if i==0 else ImageFont.truetype(FONT,16))
        notes=['Search records clearly','Use correct part number','Help next technician']
    # notes cards
    x0=55
    for n in notes:
        d.rounded_rectangle([x0,650, x0+390,725], radius=18, fill=(255,255,255), outline=(190,205,195), width=3)
        d.text((x0+20,672),n,fill=(45,65,58),font=small_f)
        x0+=430
    im.save(path)

# Keep externally generated GPT image assets if already present; otherwise create local fallback illustrations.
for fname, title, scene in [
    ('job_card_workshop.png', 'Digital Job Card in Workshop', 'jobcard'),
    ('diagnostic_scan.png', 'Diagnostic Report as a Clue', 'diagnostic'),
    ('parts_records.png', 'Parts and Service Records', 'parts'),
]:
    img_path = os.path.join(IMG_DIR, fname)
    if not os.path.exists(img_path):
        make_illustration(img_path, title, scene)

# PDF styles
styles=getSampleStyleSheet()
styles.add(ParagraphStyle('TitleMain',fontName='DejaVu-Bold',fontSize=20,leading=25,textColor=colors.HexColor('#1C5C4C'),spaceAfter=8))
styles.add(ParagraphStyle('H2Custom',fontName='DejaVu-Bold',fontSize=13,leading=17,textColor=colors.HexColor('#1C5C4C'),spaceBefore=8,spaceAfter=5))
styles.add(ParagraphStyle('BodyCustom',fontName='DejaVu',fontSize=10.2,leading=14.4,textColor=colors.HexColor('#222222'),spaceAfter=5))
styles.add(ParagraphStyle('BulletCustom',fontName='DejaVu',fontSize=10.2,leading=14,leftIndent=13,bulletIndent=4,spaceAfter=3))
styles.add(ParagraphStyle('Callout',fontName='DejaVu',fontSize=10,leading=14,textColor=colors.HexColor('#203A34'),backColor=colors.HexColor('#EEF7F1'),borderColor=colors.HexColor('#9AC7B0'),borderWidth=0.8,borderPadding=7,spaceBefore=6,spaceAfter=8))
styles.add(ParagraphStyle('Caption',fontName='DejaVu',fontSize=8.5,leading=11,textColor=colors.HexColor('#555555'),alignment=1,spaceAfter=6))


def header_footer(canvas, doc):
    w,h=A4
    if os.path.exists(HEADER):
        canvas.drawImage(HEADER, 18*mm, h-17*mm, width=174*mm, height=9.7*mm, preserveAspectRatio=True, mask='auto')
    canvas.setStrokeColor(colors.HexColor('#1C5C4C')); canvas.setLineWidth(0.7); canvas.line(18*mm,h-20*mm,192*mm,h-20*mm)
    canvas.setFont('DejaVu',7.5); canvas.setFillColor(colors.HexColor('#777777'))
    canvas.drawString(18*mm,10*mm,'© Tata Community Initiatives Trust, All Rights Reserved')
    canvas.drawRightString(192*mm,10*mm, str(doc.page))

def md_to_flow(md):
    flow=[]
    first=True
    for raw in md.split('\n'):
        line=raw.strip()
        if not line:
            flow.append(Spacer(1,2.5*mm)); continue
        if line.startswith('# '):
            if not first: flow.append(Spacer(1,2*mm))
            flow.append(Paragraph(line[2:], styles['TitleMain'])); first=False
        elif line.startswith('## '):
            flow.append(Paragraph(line[3:], styles['H2Custom']))
        elif line.startswith('- '):
            flow.append(Paragraph(line[2:], styles['BulletCustom'], bulletText='•'))
        elif line.startswith('`') and line.endswith('`'):
            flow.append(Paragraph(line, ParagraphStyle('Code',parent=styles['BodyCustom'],fontName='Courier',backColor=colors.HexColor('#F2F2F2'),borderPadding=4)))
        else:
            flow.append(Paragraph(line, styles['BodyCustom']))
    return flow

def callout(title, text):
    return Paragraph(f'<b>{title}</b><br/>{text}', styles['Callout'])

def make_rm_pdf(filename, title, md, hero, use_case_caption):
    path=os.path.join(OUT,filename)
    doc=SimpleDocTemplate(path,pagesize=A4,rightMargin=18*mm,leftMargin=18*mm,topMargin=25*mm,bottomMargin=16*mm)
    flow=[Paragraph(title, styles['TitleMain']), RLImage(hero,width=174*mm,height=94*mm), Paragraph(use_case_caption,styles['Caption'])]
    # Insert callout after intro paragraph via manual split
    for block in md_to_flow(md): flow.append(block)
    flow.append(callout('Use this in Class/Lab', 'Make a simple workshop job card for one practical task. Write the complaint, safety check, observation, action taken, and what the senior must verify.'))
    doc.build(flow,onFirstPage=header_footer,onLaterPages=header_footer)
    return path

rm1_pdf=make_rm_pdf('B2T1_RM1_EN_FORMATTED.pdf','B2T1 - Reading Material 1: What is a Digital Workspace?',rm1,os.path.join(IMG_DIR,'job_card_workshop.png'),'Figure: Learner reading a digital job card before starting practical inspection.')
rm2_pdf=make_rm_pdf('B2T1_RM2_EN_FORMATTED.pdf','B2T1 - Reading Material 2: Reading and Using Digital Information Safely',rm2,os.path.join(IMG_DIR,'diagnostic_scan.png'),'Figure: A diagnostic scan report gives clues. A mechanic still verifies practically.')

# PPT with template-like header, logo, images, footer
PPT_OUT=os.path.join(OUT,'B2T1_Learning_PPT_EN_FORMATTED.pptx')
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
PRIMARY=RGBColor(28,92,76); ACCENT=RGBColor(244,171,66); BG=RGBColor(246,248,244); DARK=RGBColor(35,45,42)

def add_header(slide, section='AI Readiness for Trades'):
    if os.path.exists(HEADER): slide.shapes.add_picture(HEADER, Inches(0.45), Inches(0.12), width=Inches(5.6))
    line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.66), Inches(13.333), Inches(0.05)); line.fill.solid(); line.fill.fore_color.rgb=PRIMARY; line.line.fill.background()
    tx=slide.shapes.add_textbox(Inches(9.4),Inches(0.12),Inches(3.4),Inches(0.35)).text_frame; tx.text=section; tx.paragraphs[0].font.size=Pt(10); tx.paragraphs[0].font.color.rgb=PRIMARY; tx.paragraphs[0].alignment=PP_ALIGN.RIGHT

def add_footer(slide, num):
    tx=slide.shapes.add_textbox(Inches(0.55),Inches(7.14),Inches(7),Inches(0.25)).text_frame; tx.text='© Tata Community Initiatives Trust, All Rights Reserved'; tx.paragraphs[0].font.size=Pt(7); tx.paragraphs[0].font.color.rgb=RGBColor(110,110,110)
    tx2=slide.shapes.add_textbox(Inches(12.45),Inches(7.14),Inches(0.4),Inches(0.25)).text_frame; tx2.text=str(num); tx2.paragraphs[0].font.size=Pt(8); tx2.paragraphs[0].font.color.rgb=RGBColor(110,110,110); tx2.paragraphs[0].alignment=PP_ALIGN.RIGHT

def add_title(slide, title, subtitle=None):
    box=slide.shapes.add_textbox(Inches(0.7),Inches(0.95),Inches(7.4),Inches(0.7)).text_frame; box.word_wrap=True; box.text=title
    p=box.paragraphs[0]; p.font.bold=True; p.font.size=Pt(31); p.font.color.rgb=DARK
    if subtitle:
        s=slide.shapes.add_textbox(Inches(0.73),Inches(1.65),Inches(7.1),Inches(0.4)).text_frame; s.text=subtitle; s.paragraphs[0].font.size=Pt(15); s.paragraphs[0].font.color.rgb=PRIMARY

def add_bullets(slide, bullets, x, y, w, h, size=18):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb=RGBColor(255,255,255); shape.line.color.rgb=RGBColor(210,222,214)
    tf=shape.text_frame; tf.margin_left=Inches(0.22); tf.margin_top=Inches(0.18); tf.word_wrap=True
    for i,b in enumerate(bullets):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=b; p.font.size=Pt(size); p.font.color.rgb=RGBColor(45,57,53); p.space_after=Pt(8)
        if i>0: p.level=0

def slide_base(num, title, subtitle=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    add_header(s); add_footer(s,num); add_title(s,title,subtitle); return s

slides=[
('B2T1 - Navigating the Digital Workspace','Mechanic Diesel and Mechanic Motor Vehicle | 50 minutes',os.path.join(IMG_DIR,'job_card_workshop.png'),['Read digital workshop information carefully','Organise records clearly','Verify before action']),
('What is a digital workspace?','It is the place where workshop information is read, updated, stored, and shared.',os.path.join(IMG_DIR,'job_card_workshop.png'),['Digital job card','Service history','Diagnostic scan report','Checklist and parts record','Photos and safety instructions']),
('Use Case 1: Digital Job Card','A learner receives a job card before starting inspection.',os.path.join(IMG_DIR,'job_card_workshop.png'),['Read the customer complaint fully','Check last service and pending work','Update only true observations','Ask senior before major action']),
('Use Case 2: Diagnostic Scan Report','A scan message is a clue, not the final repair decision.',os.path.join(IMG_DIR,'diagnostic_scan.png'),['Save or note the report','Check wiring, connector, fuse, and fitting','Use measuring tools and visual inspection','Discuss before replacing parts']),
('Use Case 3: Parts and Service Records','Good records reduce wrong parts and repeat work.',os.path.join(IMG_DIR,'parts_records.png'),['Search by job number and date','Check previous service notes','Confirm stock and correct part','Help the next technician understand the job']),
('Safe Digital Habits','Digital discipline is part of workshop discipline.',os.path.join(IMG_DIR,'diagnostic_scan.png'),['Protect customer information','Do not mark checklist items without checking','Keep devices away from oil, water, heat, and moving parts','If unsure, ask a senior']),
('Learner Workflow','A simple flow to follow in class, lab, and workplace.',os.path.join(IMG_DIR,'parts_records.png'),['Open job card → read task','Check safety instructions','Inspect practically → record observation','Get senior verification → save and close']),
('Remember','Digital tools support the mechanic. They do not replace practical skill.',os.path.join(IMG_DIR,'job_card_workshop.png'),['Read carefully','Organise clearly','Verify before action'])]
for i,(t,sub,img,bul) in enumerate(slides,1):
    s=slide_base(i,t,sub)
    s.shapes.add_picture(img, Inches(7.55), Inches(1.25), width=Inches(5.05), height=Inches(2.75))
    add_bullets(s,bul,0.75,2.35,6.35,3.25,18)
    # small orange use-case badge
    badge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.55), Inches(4.25), Inches(5.05), Inches(1.15)); badge.fill.solid(); badge.fill.fore_color.rgb=RGBColor(255,245,228); badge.line.color.rgb=ACCENT
    tf=badge.text_frame; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.13); tf.word_wrap=True
    tf.text='Workshop connection'; tf.paragraphs[0].font.bold=True; tf.paragraphs[0].font.size=Pt(15); tf.paragraphs[0].font.color.rgb=PRIMARY
    p=tf.add_paragraph(); p.text='Use digital information with hands-on inspection, safety rules, and senior guidance.'; p.font.size=Pt(12); p.font.color.rgb=RGBColor(55,55,50)
prs.save(PPT_OUT)

# Make preview images/contact sheet for QA
CONTACT=os.path.join(OUT,'B2T1_generated_images_contact_sheet.jpg')
imgs=[os.path.join(IMG_DIR,n) for n in ['job_card_workshop.png','diagnostic_scan.png','parts_records.png']]
thumbs=[Image.open(p).resize((420,228)) for p in imgs]
sheet=Image.new('RGB',(1260,280),(255,255,255)); d=ImageDraw.Draw(sheet)
for i,(p,img) in enumerate(zip(imgs,thumbs)):
    x=i*420; sheet.paste(img,(x,0)); d.text((x+10,238),os.path.basename(p),fill=(0,0,0),font=ImageFont.truetype(FONT,18))
sheet.save(CONTACT)

# Final package
zip_path=os.path.join(BASE,'B2T1_FORMATTED_RM_PPT_package.zip')
if os.path.exists(zip_path): os.remove(zip_path)
with zipfile.ZipFile(zip_path,'w',zipfile.ZIP_DEFLATED) as z:
    for root,_,files in os.walk(OUT):
        for f in files:
            p=os.path.join(root,f); z.write(p, os.path.relpath(p, BASE))
print('RM1',rm1_pdf)
print('RM2',rm2_pdf)
print('PPT',PPT_OUT)
print('ZIP',zip_path)
