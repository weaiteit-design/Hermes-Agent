from pathlib import Path
import os, zipfile, shutil, re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DInches, Pt as DPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

BASE = Path('/root/Hermes-Agent/workspace/B2T1_Navigating_Digital_Workspace')
OUT = BASE / 'template_based_final'
OUT.mkdir(exist_ok=True)
IMG = BASE / 'formatted_final' / 'images'
TEMPLATES = {
    'rm': Path('/root/.hermes/cache/documents/doc_6177a80d11b5_05_RM_Template (2).docx'),
    'intro': Path('/root/.hermes/cache/documents/doc_0dda813684d7_03_Topic Introduction Template (1).pptx'),
    'summary': Path('/root/.hermes/cache/documents/doc_2f4bcb2571d4_04_Topic Summary Template (1).pptx'),
    'video': Path('/root/.hermes/cache/documents/doc_5d585fda2fa0_06_Suggested video template (1).pptx'),
}
TOPIC = 'B2T1 - Navigating the Digital Workspace'
TOPIC_NAME = 'Navigating the Digital Workspace'

def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if hasattr(sh, 'shapes'):
            yield from iter_shapes(sh.shapes)

def delete_shape(sh):
    el = sh._element
    el.getparent().remove(el)

def set_text(shape, text, size=None, bold=None, color=None):
    if not getattr(shape, 'has_text_frame', False):
        return
    tf = shape.text_frame
    tf.clear()
    lines = str(text).split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if size or bold is not None or color:
            for run in p.runs:
                if size: run.font.size = Pt(size)
                if bold is not None: run.font.bold = bold
                if color: run.font.color.rgb = color

def replace_texts(prs, mapping):
    for slide in prs.slides:
        for sh in list(iter_shapes(slide.shapes)):
            if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
                current = '\n'.join(p.text for p in sh.text_frame.paragraphs)
                for old, new in mapping.items():
                    if old in current:
                        current = current.replace(old, new)
                if current != '\n'.join(p.text for p in sh.text_frame.paragraphs):
                    set_text(sh, current)

def remove_notes(slide):
    for sh in list(slide.shapes):
        txt = ''
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            txt = ' '.join(p.text for p in sh.text_frame.paragraphs)
        if 'Note to writer' in txt or 'DELETE the slide' in txt:
            delete_shape(sh)

def add_image_cover(slide, image_path, left, top, width, height):
    # overlay in the template-designated picture region without disturbing template elements
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

# ---------- PPT: Topic Introduction ----------
intro = Presentation(str(TEMPLATES['intro']))
# Slide 1: exact opening slide, only replace topic text
for sh in iter_shapes(intro.slides[0].shapes):
    if getattr(sh, 'has_text_frame', False) and '< Topic Number and Name >' in '\n'.join(p.text for p in sh.text_frame.paragraphs):
        set_text(sh, TOPIC, size=24, bold=True)
# Slide 2: exact MLU list layout with original groups/icons
s2 = intro.slides[1]
remove_notes(s2)
for sh in iter_shapes(s2.shapes):
    if getattr(sh, 'has_text_frame', False):
        t='\n'.join(p.text for p in sh.text_frame.paragraphs)
        if 'In this topic, you will learn about' in t:
            set_text(sh, 'In this topic, you will learn how a digital workspace helps a mechanic organise information, read records, and communicate clearly in a workshop.', size=13)
        elif '< Name of artefact >' in t:
            # fill first three artefact title placeholders in visual order
            pass
# Fill repeated artefact placeholders and descriptions in the order seen in template
placeholder_titles = ['Learning Video', 'Reading Material 1', 'Assignment']
placeholder_desc = [
    'This video will show why digital workspace skills matter in a modern automobile workshop.',
    'This reading material will explain job cards, service records, files, folders, and safe digital habits.',
    'This assignment will help you practise finding and organising workshop information.'
]
ti=di=0
for sh in iter_shapes(s2.shapes):
    if getattr(sh, 'has_text_frame', False):
        t='\n'.join(p.text for p in sh.text_frame.paragraphs)
        if '< Name of artefact >' in t and ti < len(placeholder_titles):
            set_text(sh, placeholder_titles[ti], size=13, bold=True); ti += 1
        elif 'This video will' in t or 'This reading material will' in t or 'This assignment will' in t:
            if di < len(placeholder_desc):
                set_text(sh, placeholder_desc[di], size=12); di += 1
# Slide 3: keep exact recap/assessment template slide; fit RM2 + Assessment
s3 = intro.slides[2]
remove_notes(s3)
for sh in iter_shapes(s3.shapes):
    if getattr(sh, 'has_text_frame', False):
        t='\n'.join(p.text for p in sh.text_frame.paragraphs)
        if '<topic name>' in t:
            set_text(sh, f'In this topic, you will learn about {TOPIC_NAME} with the help of the following Micro Learning Units (MLUs):', size=13)
        elif 'Topic-end Assessment' in t:
            set_text(sh, 'Reading Material 2 and Topic-end Assessment', size=13, bold=True)
# Slide 4: exact closing logo slide retained
intro.save(str(OUT / 'MMV_MD_B2T1_Topic_Introduction_TEMPLATE_FILLED.pptx'))

# ---------- PPT: Topic Summary ----------
summary = Presentation(str(TEMPLATES['summary']))
for sh in iter_shapes(summary.slides[0].shapes):
    if getattr(sh, 'has_text_frame', False) and '< Topic Name >' in '\n'.join(p.text for p in sh.text_frame.paragraphs):
        set_text(sh, TOPIC, size=24, bold=True)
# slide 2 + 3 body placeholders
recap1 = (
    'You have now reached the end of this topic.\n\n'
    'Here is what you learned:\n'
    '• A digital workspace is a place where workshop information is created, saved, found, and shared.\n'
    '• A mechanic may use job cards, service history, scan reports, part details, photos, and checklists.\n'
    '• Good file names and folders help you find information quickly during service work.\n'
    '• Always verify digital information with the instructor, service manual, machine reading, or workshop supervisor.'
)
recap2 = (
    'Some more key points:\n'
    '• A diagnostic scan report gives clues. It is not a final answer by itself.\n'
    '• Clear job notes help the next person understand what was checked and what was done.\n'
    '• Do not share customer or workshop information carelessly.\n\n'
    'My Update Habit:\n'
    '1. One new digital term/tool/output I noticed today: ________\n'
    '2. One thing I will always verify before acting: ________\n'
    '3. One reliable place I can check when unsure: ________'
)
for idx, body in [(1,recap1),(2,recap2)]:
    slide=summary.slides[idx]
    remove_notes(slide)
    for sh in iter_shapes(slide.shapes):
        if getattr(sh,'has_text_frame',False):
            t='\n'.join(p.text for p in sh.text_frame.paragraphs)
            if '<list the key concepts>' in t or 'You have now reached' in t or 'Some more key points' in t:
                set_text(sh, body, size=20 if idx==1 else 18)
# Overlay relevant GPT images into exact picture placeholder regions
add_image_cover(summary.slides[1], IMG/'job_card_workshop.png', 8229599,1376378,3824271,3824271)
add_image_cover(summary.slides[2], IMG/'diagnostic_scan.png', 8229599,1376378,3824271,3824271)
summary.save(str(OUT / 'MMV_MD_B2T1_Topic_Summary_TEMPLATE_FILLED.pptx'))

# ---------- PPT: Suggested Video Template filled as PPT only, no video rendering ----------
video = Presentation(str(TEMPLATES['video']))
video_map = {
    'The Traditional vs. The Smart Site': TOPIC_NAME,
    '< Video title >': TOPIC_NAME,
    'Description / focus area of video – This video will……..': 'This video will show how a mechanic uses digital job cards, service records, and scan reports safely in a workshop.',
    'Getting Started: How People Earn a Living': TOPIC_NAME,
    'What We Offer': 'What You Will See',
}
replace_texts(video, video_map)
# Replace lorem blocks with topic-specific content while retaining template text box positions
lorems = [
    'Digital job cards help you know the customer concern, vehicle details, and work to be done.',
    'Service history and part records help you avoid repeat mistakes and choose the correct next step.',
    'Diagnostic reports are clues. A good mechanic verifies them before repairing or replacing parts.',
    'Use folders and file names so photos, reports, and notes are easy to find later.',
    'Write clear job notes: what you checked, what you found, and what action was taken.',
    'Protect customer information. Share only with the right person in the workshop.'
]
li=0
for slide in video.slides:
    for sh in iter_shapes(slide.shapes):
        if getattr(sh,'has_text_frame',False):
            t='\n'.join(p.text for p in sh.text_frame.paragraphs)
            if 'Lorem ipsum' in t and li < len(lorems):
                set_text(sh, lorems[li], size=22); li += 1
# Add relevant images on content slides, leaving the exact first and last slides untouched except title text
add_image_cover(video.slides[2], IMG/'job_card_workshop.png', 1210000, 2400000, 5200000, 3300000)
add_image_cover(video.slides[3], IMG/'diagnostic_scan.png', 1210000, 6400000, 2000000, 1400000)
add_image_cover(video.slides[4], IMG/'parts_records.png', 12800000, 3100000, 3600000, 2400000)
video.save(str(OUT / 'MMV_MD_B2T1_Learning_Video_PPT_TEMPLATE_FILLED_NO_VIDEO.pptx'))

# ---------- Reading Materials: use DOCX template as base ----------
def clear_body(doc):
    body = doc._body._element
    for child in list(body):
        if child.tag.endswith('sectPr'):
            continue
        body.remove(child)

def shade_cell(cell, fill):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), fill)
    tcPr.append(shd)

def add_p(doc, text='', style=None, bold=False, color=None, size=None, align=None):
    p=doc.add_paragraph(style=style)
    if align: p.alignment=align
    for idx,line in enumerate(text.split('\n')):
        if idx>0: p.add_run().add_break()
        r=p.add_run(line)
        r.bold=bold
        if color: r.font.color.rgb=RGBColor(*color)
        if size: r.font.size=DPt(size)
    return p

def add_callout(doc, title, body):
    t=doc.add_table(rows=1, cols=1)
    cell=t.cell(0,0)
    shade_cell(cell, 'E8F4EA')
    p=cell.paragraphs[0]
    r=p.add_run(title+'\n'); r.bold=True; r.font.size=DPt(12); r.font.color.rgb=RGBColor(39,98,54)
    r=p.add_run(body); r.font.size=DPt(10)
    return t

def build_rm(filename, title, sections, image_name):
    doc=Document(str(TEMPLATES['rm']))
    clear_body(doc)
    add_p(doc, title, 'Heading 1', bold=True, color=(39,98,54), size=18)
    doc.add_picture(str(IMG/image_name), width=DInches(5.9))
    if doc.paragraphs[-1]: doc.paragraphs[-1].alignment=WD_ALIGN_PARAGRAPH.CENTER
    add_p(doc, 'The contents of this document help learners understand the topic in simple language and connect it with class, lab, and workshop situations.', size=10)
    for heading, body in sections:
        if heading.startswith('Think') or heading.startswith('Use') or heading.startswith('Key'):
            add_callout(doc, heading, body)
        else:
            add_p(doc, heading, 'Heading 2', bold=True, color=(39,98,54), size=14)
            for para in body.split('\n\n'):
                add_p(doc, para, 'normal', size=10)
    doc.save(str(OUT/filename))

rm1_sections=[
('What is a Digital Workspace?', 'A digital workspace is the place where work information is created, saved, found, and shared using a computer, tablet, mobile device, or workshop system.\n\nIn a vehicle workshop, it may include job cards, service history, inspection photos, diagnostic scan reports, parts details, checklists, and messages from the supervisor.'),
('Why it matters in a workshop', 'A mechanic works with tools, machines, and information. If the information is wrong or missing, the repair can become slow or unsafe.\n\nExample: A learner sees a job card that says “engine starting problem”. Before opening parts, the learner checks the complaint, previous service note, battery reading, and instructor instruction.'),
('Use this in Class/Lab', 'Open a sample job card. Identify: customer complaint, vehicle details, date, assigned work, and safety note. Then explain what you will check first.'),
('Good digital habits', '• Name files clearly, such as jobcard_engine_starting_problem.\n• Keep photos and reports in the correct folder.\n• Do not delete workshop records without permission.\n• Verify information before using it for repair decisions.\n• Do not share customer information outside the workshop.'),
('Key Takeaways', 'Digital workspace means organised digital information for work. A mechanic should know how to find, read, save, and verify workshop information.')]
rm2_sections=[
('Reading digital workshop information', 'A diagnostic scan report, service history, or part record can support a mechanic. But digital output is not a final answer by itself. It gives clues. The mechanic must verify the clue with observation, measurement, manual, and instructor guidance.'),
('Use Case 1: Diagnostic report as a clue', 'A scan report shows a sensor-related fault. A learner should not replace the sensor immediately. First check wiring, connector condition, battery voltage, and whether the fault returns after clearing. This saves time and avoids unnecessary replacement.'),
('Use Case 2: Finding the right service record', 'A vehicle returns with the same complaint. The mechanic checks previous job notes and sees that a related check was already done. This helps the workshop avoid repeating the same work and plan the next check properly.'),
('Use Case 3: Parts information', 'A part may look similar but have different specifications. The learner checks the part record and confirms with the instructor before selecting it.'),
('Think & Reflect', 'What can go wrong if a mechanic follows a digital report without checking the actual vehicle? Write two points.'),
('Key Takeaways', 'Digital information helps only when it is organised, understood, and verified. A good mechanic combines digital clues with practical checking.')]
build_rm('MMV_MD_B2T1_RM_A1_TEMPLATE_FILLED.docx', 'Reading Material 1: Digital Workspace Basics', rm1_sections, 'job_card_workshop.png')
build_rm('MMV_MD_B2T1_RM_A2_TEMPLATE_FILLED.docx', 'Reading Material 2: Using Digital Information Safely', rm2_sections, 'diagnostic_scan.png')

# ---------- Package ----------
zpath=BASE/'B2T1_TEMPLATE_FOLLOWED_RM_PPT_package.zip'
if zpath.exists(): zpath.unlink()
with zipfile.ZipFile(zpath,'w',zipfile.ZIP_DEFLATED) as z:
    for p in OUT.rglob('*'):
        if p.is_file(): z.write(p, p.relative_to(BASE))
print(zpath)
print('files', len(list(OUT.rglob('*'))))
